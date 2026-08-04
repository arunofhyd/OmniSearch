#!/usr/bin/env python3
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette
    COLOR_BG = RGBColor(0, 0, 0)
    COLOR_CARD = RGBColor(17, 17, 17)
    COLOR_CARD_BORDER = RGBColor(40, 40, 40)
    COLOR_WHITE = RGBColor(255, 255, 255)
    COLOR_SILVER = RGBColor(180, 180, 185)
    COLOR_GRAY = RGBColor(134, 134, 139)
    COLOR_BLUE = RGBColor(41, 151, 255)
    COLOR_PURPLE = RGBColor(191, 90, 242)
    COLOR_GREEN = RGBColor(48, 209, 88)

    blank_layout = prs.slide_layouts[6]

    def set_slide_bg(slide):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_BG

    def add_header(slide, tag_text, title_text, tag_color=COLOR_BLUE):
        if tag_text:
            tb_tag = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.4))
            tf_tag = tb_tag.text_frame
            tf_tag.word_wrap = True
            p_tag = tf_tag.paragraphs[0]
            p_tag.text = tag_text.upper()
            p_tag.font.name = "Arial"
            p_tag.font.size = Pt(11)
            p_tag.font.bold = True
            p_tag.font.color.rgb = tag_color

        if title_text:
            tb_title = slide.shapes.add_textbox(Inches(0.8), Inches(0.9), Inches(11.7), Inches(0.8))
            tf_title = tb_title.text_frame
            tf_title.word_wrap = True
            p_title = tf_title.paragraphs[0]
            p_title.text = title_text
            p_title.font.name = "Arial"
            p_title.font.size = Pt(32)
            p_title.font.bold = True
            p_title.font.color.rgb = COLOR_WHITE

    # -------------------------------------------------------------------------
    # SLIDE 1: TITLE
    # -------------------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide1)

    tb = slide1.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(11.3), Inches(2.0))
    tf = tb.text_frame
    p1 = tf.paragraphs[0]
    p1.text = "Omni Search"
    p1.font.name = "Arial"
    p1.font.size = Pt(64)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_WHITE
    p1.alignment = PP_ALIGN.LEFT

    p2 = tf.add_paragraph()
    p2.text = "Search at the speed of thought."
    p2.font.name = "Arial"
    p2.font.size = Pt(28)
    p2.font.color.rgb = COLOR_BLUE
    p2.alignment = PP_ALIGN.LEFT

    # Presenter Pill
    tb_pres = slide1.shapes.add_textbox(Inches(9.0), Inches(6.2), Inches(3.5), Inches(0.6))
    tf_pres = tb_pres.text_frame
    p_pres = tf_pres.paragraphs[0]
    p_pres.text = "Presented by Arun Thomas"
    p_pres.font.name = "Arial"
    p_pres.font.size = Pt(14)
    p_pres.font.color.rgb = COLOR_GRAY
    p_pres.alignment = PP_ALIGN.RIGHT

    # -------------------------------------------------------------------------
    # SLIDE 2: PROBLEM STATEMENT
    # -------------------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide2)
    add_header(slide2, "01 Problem Statement", "The daily struggle is real.")

    cards_data2 = [
        ("Context-Switching Tax", "Every search requires leaving the current app, opening a browser, navigating to the right tool, typing the query, and switching back. Analysts lose 5–15 minutes per day to this invisible friction."),
        ("Tab Overload", "No standardized search workflow. Each analyst uses different tools and methods, leading to browser clutter, lost tabs, and fragmented knowledge across dozens of open windows."),
        ("Accessibility Gap", "Specially-abled and blind analysts face an amplified burden — multi-step, mouse-dependent search workflows are a barrier. A keyboard-first, VoiceOver-compatible solution is essential.")
    ]

    for i, (ctitle, cdesc) in enumerate(cards_data2):
        left = Inches(0.8 + i * 3.9)
        top = Inches(2.0)
        width = Inches(3.7)
        height = Inches(4.5)

        shape = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_CARD
        shape.line.color.rgb = COLOR_CARD_BORDER

        tf = shape.text_frame
        tf.margin_left = Inches(0.3)
        tf.margin_top = Inches(0.4)
        tf.margin_right = Inches(0.3)
        tf.word_wrap = True

        p_t = tf.paragraphs[0]
        p_t.text = ctitle
        p_t.font.size = Pt(20)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_WHITE
        p_t.space_after = Pt(14)

        p_d = tf.add_paragraph()
        p_d.text = cdesc
        p_d.font.size = Pt(14)
        p_d.font.color.rgb = COLOR_GRAY

    # -------------------------------------------------------------------------
    # SLIDE 3: OBJECTIVE
    # -------------------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide3)
    add_header(slide3, "02 Objective", "Native, Zero-Friction & Accessible")

    tb3 = slide3.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.33), Inches(1.8))
    tf3 = tb3.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    p3.text = "Build a native, zero-friction search utility that routes any selected text to the right tool — instantly — from any app, on any desktop, for every analyst."
    p3.font.size = Pt(26)
    p3.font.color.rgb = COLOR_WHITE
    p3.alignment = PP_ALIGN.LEFT

    pillars = [
        ("Zero Friction", COLOR_BLUE),
        ("Fully Accessible", COLOR_GREEN),
        ("Intelligent Routing", COLOR_PURPLE)
    ]
    for i, (pname, pcolor) in enumerate(pillars):
        left = Inches(1.0 + i * 3.9)
        top = Inches(4.5)
        width = Inches(3.5)
        height = Inches(1.8)

        shape = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_CARD
        shape.line.color.rgb = pcolor

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = pname
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = pcolor
        p.alignment = PP_ALIGN.CENTER

    # -------------------------------------------------------------------------
    # SLIDE 4: SOLUTION — HERO CARD
    # -------------------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide4)
    add_header(slide4, "03 Solution", "Omni Search. Absolute focus.")

    tb4_desc = slide4.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(0.8))
    tf4_desc = tb4_desc.text_frame
    tf4_desc.word_wrap = True
    p4_d = tf4_desc.paragraphs[0]
    p4_d.text = "A native macOS shortcut utility that routes queries directly to your favorite tools, cleans search URLs automatically, and reuses open browser windows without distracting tab clutter."
    p4_d.font.size = Pt(16)
    p4_d.font.color.rgb = COLOR_GRAY

    feats4 = [
        ("Lightning Fast", "Reuses active windows instantly via precise window targeting."),
        ("Smart Link Sanitizer", "Cleans and formats search URLs so they open smoothly in enterprise tools without security block errors."),
        ("Instant Focus", "Pulls search results to the absolute front across all virtual desktops.")
    ]

    for i, (ftitle, fdesc) in enumerate(feats4):
        left = Inches(0.8 + i * 3.9)
        top = Inches(2.8)
        width = Inches(3.7)
        height = Inches(3.8)

        shape = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_CARD
        shape.line.color.rgb = COLOR_CARD_BORDER

        tf = shape.text_frame
        tf.margin_left = Inches(0.3)
        tf.margin_top = Inches(0.4)
        tf.word_wrap = True

        p_t = tf.paragraphs[0]
        p_t.text = ftitle
        p_t.font.size = Pt(18)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_BLUE
        p_t.space_after = Pt(10)

        p_d = tf.add_paragraph()
        p_d.text = fdesc
        p_d.font.size = Pt(14)
        p_d.font.color.rgb = COLOR_GRAY

    # -------------------------------------------------------------------------
    # SLIDE 5: HOW IT WORKS
    # -------------------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide5)
    add_header(slide5, "03 Solution", "How does it work?")

    steps = [
        ("1", "Highlight", "Select any text, in any app, anywhere on your screen.", COLOR_BLUE),
        ("2", "Trigger", "Execute via global keyboard hotkey or right-click mouse.", COLOR_PURPLE),
        ("3", "Route", "Route the relevant search target to get your result.", COLOR_GREEN)
    ]

    for i, (num, stitle, sdesc, scolor) in enumerate(steps):
        left = Inches(0.8 + i * 3.9)
        top = Inches(2.2)
        width = Inches(3.7)
        height = Inches(4.2)

        shape = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_CARD
        shape.line.color.rgb = COLOR_CARD_BORDER

        tf = shape.text_frame
        tf.margin_left = Inches(0.3)
        tf.margin_top = Inches(0.4)
        tf.word_wrap = True

        p_num = tf.paragraphs[0]
        p_num.text = num
        p_num.font.size = Pt(36)
        p_num.font.bold = True
        p_num.font.color.rgb = scolor
        p_num.space_after = Pt(10)

        p_t = tf.add_paragraph()
        p_t.text = stitle
        p_t.font.size = Pt(22)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_WHITE
        p_t.space_after = Pt(10)

        p_d = tf.add_paragraph()
        p_d.text = sdesc
        p_d.font.size = Pt(14)
        p_d.font.color.rgb = COLOR_GRAY

    # -------------------------------------------------------------------------
    # SLIDE 6: CUSTOMIZATION & LOCALES
    # -------------------------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide6)
    add_header(slide6, "03 Solution", "Customization & Global Locales")

    tb6 = slide6.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tf6 = tb6.text_frame
    tf6.word_wrap = True

    items6 = [
        ("Locale & Region Friendly", "Search engines and tools support any locale seamlessly: en_IN, en_US, en_GB, en_SA, + any region."),
        ("Tailored Workflows", "Customize menu options for Apple Music, Marketing Media Tools, Google Search, and internal databases."),
        ("Centralized Hub", "Consolidate scattered bookmarks and enterprise tools into one quick menu.")
    ]

    for title, desc in items6:
        p_t = tf6.add_paragraph() if tf6.paragraphs[0].text else tf6.paragraphs[0]
        p_t.text = f"•  {title}"
        p_t.font.size = Pt(20)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_BLUE

        p_d = tf6.add_paragraph()
        p_d.text = f"    {desc}"
        p_d.font.size = Pt(15)
        p_d.font.color.rgb = COLOR_GRAY
        p_d.space_after = Pt(18)

    # -------------------------------------------------------------------------
    # SLIDE 7: TIME SAVED
    # -------------------------------------------------------------------------
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide7)
    add_header(slide7, "04 Time Saved", "The numbers speak.")

    kpis = [
        ("8.5 Hours", "Saved per analyst / month"),
        ("94% Reduction", "Search time cut from 30s → 2s"),
        ("100% Native", "Zero dependencies, fully accessible")
    ]

    for i, (kval, klabel) in enumerate(kpis):
        left = Inches(0.8 + i * 3.9)
        top = Inches(2.2)
        width = Inches(3.7)
        height = Inches(4.2)

        shape = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_CARD
        shape.line.color.rgb = COLOR_CARD_BORDER

        tf = shape.text_frame
        tf.margin_left = Inches(0.3)
        tf.margin_top = Inches(0.5)
        tf.word_wrap = True

        p_v = tf.paragraphs[0]
        p_v.text = kval
        p_v.font.size = Pt(28)
        p_v.font.bold = True
        p_v.font.color.rgb = COLOR_GREEN
        p_v.space_after = Pt(14)

        p_l = tf.add_paragraph()
        p_l.text = klabel
        p_l.font.size = Pt(15)
        p_l.font.color.rgb = COLOR_GRAY

    # -------------------------------------------------------------------------
    # SLIDE 8: PILOT FEEDBACK
    # -------------------------------------------------------------------------
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide8)
    add_header(slide8, "05 Pilot Feedback", "What analysts are saying.")

    quotes = [
        ('"Saved me at least 15 minutes of tab-hunting every day. Indispensable for high-volume research."', "Senior Data Analyst", "Media Research Team"),
        ('"VoiceOver navigation is flawless. First tool that treats accessibility as a core feature."', "Accessibility Specialist", "Product Operations"),
        ('"Sanitizing URL parameters automatically saved our team from broken enterprise tool links."', "Operations Lead", "Workflow Engineering")
    ]

    for i, (qtext, qauthor, qrole) in enumerate(quotes):
        left = Inches(0.8 + i * 3.9)
        top = Inches(2.2)
        width = Inches(3.7)
        height = Inches(4.2)

        shape = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_CARD
        shape.line.color.rgb = COLOR_CARD_BORDER

        tf = shape.text_frame
        tf.margin_left = Inches(0.3)
        tf.margin_top = Inches(0.4)
        tf.word_wrap = True

        p_q = tf.paragraphs[0]
        p_q.text = qtext
        p_q.font.size = Pt(14)
        p_q.font.italic = True
        p_q.font.color.rgb = COLOR_WHITE
        p_q.space_after = Pt(20)

        p_a = tf.add_paragraph()
        p_a.text = qauthor
        p_a.font.size = Pt(14)
        p_a.font.bold = True
        p_a.font.color.rgb = COLOR_BLUE

        p_r = tf.add_paragraph()
        p_r.text = qrole
        p_r.font.size = Pt(12)
        p_r.font.color.rgb = COLOR_GRAY

    # -------------------------------------------------------------------------
    # SLIDE 9: MIC DROP / CTA
    # -------------------------------------------------------------------------
    slide9 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide9)

    tb9 = slide9.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(11.3), Inches(3.0))
    tf9 = tb9.text_frame
    tf9.word_wrap = True

    p9_1 = tf9.paragraphs[0]
    p9_1.text = "Fluid. Inclusive. Native. Universal."
    p9_1.font.size = Pt(44)
    p9_1.font.bold = True
    p9_1.font.color.rgb = COLOR_WHITE

    p9_2 = tf9.add_paragraph()
    p9_2.text = "Ready for immediate deployment across analyst workstations."
    p9_2.font.size = Pt(22)
    p9_2.font.color.rgb = COLOR_BLUE
    p9_2.space_before = Pt(16)

    # -------------------------------------------------------------------------
    # SLIDE 10: DOWNLOAD & RELEASES
    # -------------------------------------------------------------------------
    slide10 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide10)
    add_header(slide10, "Releases", "Select Your Version")

    versions = [
        ("v1.1 (Latest Release)", "With Auto-Updater", "Silent background update checks, link sanitizer, and PID tracking."),
        ("v1.0 (Stable Release)", "Safe & Production Ready", "Original Omni Search native experience with manual updates.")
    ]

    for i, (vnum, vsub, vdesc) in enumerate(versions):
        left = Inches(1.5 + i * 5.2)
        top = Inches(2.2)
        width = Inches(4.8)
        height = Inches(4.2)

        shape = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_CARD
        shape.line.color.rgb = COLOR_BLUE if i == 0 else COLOR_CARD_BORDER

        tf = shape.text_frame
        tf.margin_left = Inches(0.4)
        tf.margin_top = Inches(0.5)
        tf.word_wrap = True

        p_v = tf.paragraphs[0]
        p_v.text = vnum
        p_v.font.size = Pt(24)
        p_v.font.bold = True
        p_v.font.color.rgb = COLOR_WHITE
        p_v.space_after = Pt(6)

        p_s = tf.add_paragraph()
        p_s.text = vsub
        p_s.font.size = Pt(16)
        p_s.font.color.rgb = COLOR_BLUE

        p_d = tf.add_paragraph()
        p_d.text = vdesc
        p_d.font.size = Pt(14)
        p_d.font.color.rgb = COLOR_GRAY
        p_d.space_before = Pt(14)

    output_path = "/Users/arunthomas/.gemini/antigravity-ide/scratch/OmniSearch/OmniSearch_Presentation.pptx"
    prs.save(output_path)
    print(f"SUCCESS: Generated presentation at {output_path}")

if __name__ == "__main__":
    create_presentation()
